"""Generate the WeatherLens PowerPoint presentation.

Output: ``WeatherLens_Prezentare.pptx`` (project root) — 18 slides, structured for
the academic presentation:

  1.  Title
  2.  What is AWS?
  3.  AWS service categories
  4.  WeatherLens — the application
  5.  Project scope & goals
  6.  The 9 services (overview list)
  7-15. Each AWS service detailed (one slide per service)
  16. Connection schema (how the 9 services link together)
  17. Macro workflow (4 phases — Ingest → Process → Store → Display)
  18. Conclusion / Thank you

Re-run with ``python scripts/build_presentation.py``. Idempotent.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUTPUT = ROOT / "WeatherLens_Prezentare.pptx"

# ── Theme ────────────────────────────────────────────────────────────
BG_DARK    = RGBColor(0x0F, 0x17, 0x2A)
BG_SURFACE = RGBColor(0x1E, 0x29, 0x3B)
BG_DEEP    = RGBColor(0x0B, 0x12, 0x20)
ACCENT     = RGBColor(0x3B, 0x82, 0xF6)   # blue
ACCENT_HOT = RGBColor(0xEF, 0x44, 0x44)   # red
ACCENT_GR  = RGBColor(0x10, 0xB9, 0x81)   # green
ACCENT_AM  = RGBColor(0xF5, 0x9E, 0x0B)   # amber
ACCENT_VL  = RGBColor(0xA7, 0x8B, 0xFA)   # violet
ACCENT_CY  = RGBColor(0x06, 0xB6, 0xD4)   # cyan
TEXT_LIGHT = RGBColor(0xE2, 0xE8, 0xF0)
TEXT_MUTED = RGBColor(0x94, 0xA3, 0xB8)
TEXT_DIM   = RGBColor(0x64, 0x74, 0x8B)

WIDE_W = Inches(13.333)
WIDE_H = Inches(7.5)


# ── Primitives ──────────────────────────────────────────────────────

def _set_bg(slide, color=BG_DARK):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, WIDE_W, WIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def _add_text(slide, x, y, w, h, text, *, size=18, bold=False,
              color=TEXT_LIGHT, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return tf


def _add_paragraphs(slide, x, y, w, h, lines, *, size=14,
                    color=TEXT_LIGHT, align=PP_ALIGN.LEFT, space_after=4):
    """Like _add_text but accepts a list of lines (each becomes a paragraph)."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.space_after = Pt(space_after)
        run = para.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return tf


def _add_bullets(slide, x, y, w, h, items, *, size=14, color=TEXT_LIGHT,
                 space_after=6, bullet_char="•"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        para.space_after = Pt(space_after)
        run = para.add_run()
        run.text = f"{bullet_char}  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return tf


def _add_card(slide, x, y, w, h, *, fill=BG_SURFACE, border=ACCENT,
              border_width=1.0):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = border
    card.line.width = Pt(border_width)
    card.adjustments[0] = 0.06
    card.shadow.inherit = False
    return card


def _add_accent_bar(slide, color=ACCENT):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, WIDE_W, Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def _slide_title(slide, title, subtitle=None, accent=ACCENT):
    _add_accent_bar(slide, accent)
    _add_text(slide, Inches(0.5), Inches(0.30), Inches(12), Inches(0.85),
              title, size=32, bold=True, color=TEXT_LIGHT)
    if subtitle:
        _add_text(slide, Inches(0.5), Inches(1.05), Inches(12), Inches(0.5),
                  subtitle, size=15, color=TEXT_MUTED)


def _new_blank(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _set_bg(slide)
    return slide


def _twin_callouts(slide, *, left_label, left_text, right_label, right_text,
                   right_color=ACCENT_GR, top=Inches(1.70)):
    """Two side-by-side callout cards.

    Left = generic ('Despre serviciu') with violet accent.
    Right = project-specific ('Rol in WeatherLens') with the service's color.

    Both fixed height; caller can place content below at top + 1.55in.
    """
    h = Inches(1.55)
    half_w = Inches(6.05)
    gap = Inches(0.20)
    x_left = Inches(0.5)
    x_right = x_left + half_w + gap

    # Left card
    _add_card(slide, x_left, top, half_w, h, fill=BG_SURFACE, border=ACCENT_VL,
              border_width=1.2)
    stripe_l = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      x_left, top + Inches(0.08),
                                      Inches(0.10), h - Inches(0.16))
    stripe_l.fill.solid()
    stripe_l.fill.fore_color.rgb = ACCENT_VL
    stripe_l.line.fill.background()
    stripe_l.shadow.inherit = False
    _add_text(slide, x_left + Inches(0.30), top + Inches(0.12),
              half_w - Inches(0.4), Inches(0.32),
              left_label, size=10, bold=True, color=ACCENT_VL)
    _add_text(slide, x_left + Inches(0.30), top + Inches(0.42),
              half_w - Inches(0.4), h - Inches(0.5),
              left_text, size=12, color=TEXT_LIGHT)

    # Right card
    _add_card(slide, x_right, top, half_w, h, fill=BG_SURFACE, border=right_color,
              border_width=1.2)
    stripe_r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      x_right, top + Inches(0.08),
                                      Inches(0.10), h - Inches(0.16))
    stripe_r.fill.solid()
    stripe_r.fill.fore_color.rgb = right_color
    stripe_r.line.fill.background()
    stripe_r.shadow.inherit = False
    _add_text(slide, x_right + Inches(0.30), top + Inches(0.12),
              half_w - Inches(0.4), Inches(0.32),
              right_label, size=10, bold=True, color=right_color)
    _add_text(slide, x_right + Inches(0.30), top + Inches(0.42),
              half_w - Inches(0.4), h - Inches(0.5),
              right_text, size=12, color=TEXT_LIGHT)

    return top + h  # bottom Y so caller can place content below


# ════════════════════════════════════════════════════════════════════
# SECTION 1 — TITLE
# ════════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    s = _new_blank(prs)
    # Subtle accent block
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(0), Inches(2.15), WIDE_W, Inches(0.04))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()
    accent.shadow.inherit = False

    _add_text(s, Inches(0.5), Inches(2.4), Inches(12.3), Inches(1.0),
              "☁️  WeatherLens",
              size=68, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
    _add_text(s, Inches(0.5), Inches(3.65), Inches(12.3), Inches(0.6),
              "AWS Data Analytics Platform",
              size=28, color=ACCENT, align=PP_ALIGN.CENTER)
    _add_text(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.5),
              "Demonstrarea serviciilor AWS prin construirea unei aplicatii reale",
              size=18, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    _add_text(s, Inches(0.5), Inches(5.75), Inches(12.3), Inches(0.5),
              "Proiect IOC · AlexOnData · 2026",
              size=14, color=TEXT_DIM, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SECTION 2 — AWS INTRO
# ════════════════════════════════════════════════════════════════════

def slide_02_what_is_aws(prs):
    s = _new_blank(prs)
    _slide_title(s, "Ce este AWS?",
                 "Amazon Web Services — cea mai mare platforma cloud din lume")

    # Big stats row
    stats = [
        ("200+",        "servicii cloud",          ACCENT),
        ("32",          "regiuni globale",         ACCENT_GR),
        ("#1",          "cota de piata cloud",     ACCENT_AM),
        ("$100B+",      "venit anual (2024)",      ACCENT_HOT),
    ]
    card_w = Inches(2.85)
    card_h = Inches(1.55)
    gap = Inches(0.25)
    total_w = 4 * card_w + 3 * gap
    start_x = (WIDE_W - total_w) / 2
    y = Inches(2.0)
    for i, (number, label, color) in enumerate(stats):
        x = start_x + i * (card_w + gap)
        _add_card(s, x, y, card_w, card_h, fill=BG_SURFACE, border=color)
        _add_text(s, x, y + Inches(0.25), card_w, Inches(0.6),
                  number, size=36, bold=True, color=color, align=PP_ALIGN.CENTER)
        _add_text(s, x, y + Inches(1.00), card_w, Inches(0.4),
                  label, size=13, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    # Description
    bullets = [
        "Platforma cloud lansata de Amazon in 2006 — pionierul modelului infrastructure-as-a-service.",
        "Oferta acopera tot ce ai nevoie: storage, compute, baze de date, ML, securitate, retele, BI.",
        "Modelul de cost: pay-as-you-go — platesti exact ce folosesti, la secunda sau per request.",
        "Folosit de Netflix, Airbnb, NASA, BMW, Shell — scaleaza de la startup la giganti.",
    ]
    _add_bullets(s, Inches(0.7), Inches(4.0), Inches(12), Inches(2.8),
                 bullets, size=14, space_after=8)


def slide_03_aws_categories(prs):
    s = _new_blank(prs)
    _slide_title(s, "Categorii de servicii AWS",
                 "200+ servicii organizate in 8 mari categorii — WeatherLens foloseste 6")

    categories = [
        ("Storage",       "S3 · EBS · Glacier",            ACCENT,     True),
        ("Compute",       "EC2 · Lambda · Fargate",        ACCENT_AM,  True),
        ("Database",      "RDS · DynamoDB · Aurora",       ACCENT_VL,  False),
        ("Networking",    "VPC · Route 53 · CloudFront",   ACCENT_CY,  False),
        ("Analytics",     "Athena · Glue · EMR · Kinesis", ACCENT_HOT, True),
        ("Machine Learning", "SageMaker · Rekognition · Bedrock", ACCENT_VL, False),
        ("Security & Ops","IAM · CloudWatch · KMS · WAF",  ACCENT_GR,  True),
        ("Integration",   "EventBridge · SNS · SQS · Step Functions", ACCENT, True),
    ]
    cols = 4
    cell_w = Inches(3.05)
    cell_h = Inches(1.55)
    gap_x = Inches(0.18)
    gap_y = Inches(0.18)
    start_x = Inches(0.4)
    start_y = Inches(1.95)
    for i, (name, examples, color, used) in enumerate(categories):
        row, col = divmod(i, cols)
        x = start_x + col * (cell_w + gap_x)
        y = start_y + row * (cell_h + gap_y)
        fill = BG_SURFACE if used else BG_DEEP
        bw = 2.0 if used else 0.8
        _add_card(s, x, y, cell_w, cell_h, fill=fill, border=color, border_width=bw)
        _add_text(s, x + Inches(0.2), y + Inches(0.2), cell_w - Inches(0.4),
                  Inches(0.45), name, size=16, bold=True, color=color)
        _add_text(s, x + Inches(0.2), y + Inches(0.7), cell_w - Inches(0.4),
                  Inches(0.6), examples, size=11, color=TEXT_LIGHT)
        if used:
            _add_text(s, x + Inches(0.2), y + Inches(1.15), cell_w - Inches(0.4),
                      Inches(0.3), "✓ folosit in WeatherLens",
                      size=10, bold=True, color=ACCENT_GR)

    _add_text(s, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4),
              "Categoriile cu chenar puternic = folosite in proiect · 6 din 8 categorii acoperite",
              size=11, color=TEXT_MUTED, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SECTION 3 — WEATHERLENS APP & SCOPE
# ════════════════════════════════════════════════════════════════════

def slide_04_weatherlens_app(prs):
    s = _new_blank(prs)
    _slide_title(s, "WeatherLens — aplicatia",
                 "Pipeline meteorologic complet construit pe AWS")

    # Hero description card
    _add_card(s, Inches(0.5), Inches(1.75), Inches(12.3), Inches(1.4),
              fill=BG_SURFACE, border=ACCENT, border_width=1.5)
    _add_text(s, Inches(0.8), Inches(1.95), Inches(11.7), Inches(0.45),
              "Ce face aplicatia",
              size=14, bold=True, color=ACCENT)
    _add_text(s, Inches(0.8), Inches(2.32), Inches(11.7), Inches(0.85),
              "WeatherLens preia automat date meteo orare de la Open-Meteo API pentru "
              "5 orase romanesti, le proceseaza printr-un pipeline ETL serverless si le "
              "expune intr-un dashboard BI cu grafice interactive si export Excel/CSV.",
              size=14, color=TEXT_LIGHT)

    # Three feature tiles
    features = [
        ("📡 Ingestie automata",
         "Cron zilnic preia 24 de inregistrari × 5 orase = 120 randuri/zi de la Open-Meteo. "
         "Fara interventie manuala dupa deploy."),
        ("⚙️ ETL serverless",
         "JSON → Parquet partitionat cu 13 derivari analitice (sezon, anomalii, "
         "rolling avg, categorii vant, fenomene extreme). Validare ranges automat."),
        ("📊 Dashboard interactiv",
         "Filtre dinamice (oras, perioada, sezon), 7 tipuri de grafice, KPI cu delta "
         "vs perioada precedenta, export CSV/Excel/JSON."),
    ]
    y = Inches(3.40)
    for i, (title, desc) in enumerate(features):
        h = Inches(1.05)
        _add_card(s, Inches(0.5), y, Inches(12.3), h, fill=BG_SURFACE, border=ACCENT_GR)
        _add_text(s, Inches(0.8), y + Inches(0.13), Inches(11.7), Inches(0.4),
                  title, size=15, bold=True, color=ACCENT_GR)
        _add_text(s, Inches(0.8), y + Inches(0.48), Inches(11.7), Inches(0.55),
                  desc, size=12, color=TEXT_LIGHT)
        y += h + Inches(0.13)


def slide_05_scope(prs):
    s = _new_blank(prs)
    _slide_title(s, "Scopul proiectului",
                 "WeatherLens = vehicul didactic pentru explorarea ecosistemului AWS")

    # Two-column layout
    _add_card(s, Inches(0.5), Inches(1.75), Inches(6.0), Inches(4.6),
              fill=BG_SURFACE, border=ACCENT)
    _add_text(s, Inches(0.7), Inches(1.95), Inches(5.7), Inches(0.45),
              "OBIECTIV ACADEMIC",
              size=12, bold=True, color=ACCENT)
    _add_text(s, Inches(0.7), Inches(2.35), Inches(5.7), Inches(3.95),
              "Prezentarea unei arhitecturi cloud reale — nu un demo teoretic, "
              "ci un sistem care:\n\n"
              "• Acopera toate cele 6 layere ale unei platforme de date\n"
              "  (Trigger · Ingestion · Orchestration · Processing · Analytics · Visualization)\n\n"
              "• Demonstreaza integrare reala intre 9 servicii AWS\n"
              "  — fiecare cu un rol specific in pipeline\n\n"
              "• Include un serviciu nou pentru cerinta de Nota 10\n"
              "  — Step Functions ⭐",
              size=13, color=TEXT_LIGHT)

    _add_card(s, Inches(6.8), Inches(1.75), Inches(6.0), Inches(4.6),
              fill=BG_SURFACE, border=ACCENT_GR)
    _add_text(s, Inches(7.0), Inches(1.95), Inches(5.7), Inches(0.45),
              "REZULTATE LIVRATE",
              size=12, bold=True, color=ACCENT_GR)
    _add_text(s, Inches(7.0), Inches(2.35), Inches(5.7), Inches(3.95),
              "Toate cerintele bifate, plus extra:\n\n"
              "• 9 servicii AWS integrate (cerinta minim 8)\n\n"
              "• Pipeline complet automatizat zilnic\n"
              "  ETL ruleaza fara supraveghere la 06:00 UTC\n\n"
              "• Dashboard 4 pagini cu drill-down si export\n\n"
              "• Bonus: implementare locala $0 cu mapping AWS↔local\n"
              "  (Streamlit + DuckDB + Python)\n\n"
              "• 12 teste pytest verzi · documentatie completa",
              size=13, color=TEXT_LIGHT)

    # Footer KPI
    _add_card(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.6),
              fill=BG_SURFACE, border=ACCENT_HOT)
    _add_text(s, Inches(0.7), Inches(6.66), Inches(12.0), Inches(0.4),
              "Date reale: 5 orase · 119 zile (2026-01-01 → 2026-04-29) · 14 280 randuri orare · 595 rezumate zilnice",
              size=12, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SECTION 4 — SERVICES OVERVIEW (one-sentence each)
# ════════════════════════════════════════════════════════════════════

def slide_06_services_list(prs):
    s = _new_blank(prs)
    _slide_title(s, "Cele 9 servicii AWS folosite",
                 "Lista completa cu rolul fiecaruia in WeatherLens")

    services = [
        ("📦", "Amazon S3",          ACCENT,
         "Stocheaza datele brute (JSON), procesate (Parquet), exporturile "
         "si site-ul static al frontend-ului."),
        ("🔧", "AWS Lambda",         ACCENT_AM,
         "Ruleaza cele 3 functii Python ale proiectului (ingest · validate · export) "
         "fara servere de mentenut."),
        ("⚙️", "AWS Glue",           ACCENT_HOT,
         "Motorul ETL care transforma JSON → Parquet partitionat cu 13 derivari "
         "analitice in PySpark."),
        ("🎯", "AWS Step Functions", ACCENT_GR,
         "Orchestreaza pipeline-ul ETL printr-un state machine vizual cu retry "
         "exponential si catch automat."),
        ("🗄️", "Amazon Athena",      ACCENT,
         "Motor SQL serverless care ruleaza cele 10 query-uri pre-definite "
         "direct pe Parquet din S3."),
        ("📊", "Amazon QuickSight",  ACCENT_HOT,
         "Dashboard BI cu 7 grafice interactive, filtre dinamice si refresh "
         "automat dupa pipeline."),
        ("📈", "Amazon CloudWatch",  ACCENT_AM,
         "Centralizeaza logs si metrici de la toate serviciile, plus alarme "
         "email cand pipeline-ul esueaza."),
        ("🔐", "AWS IAM",            ACCENT_GR,
         "Defineste 3 roluri least-privilege pentru Glue, Lambda si Step Functions; "
         "zero credentials hardcoded."),
        ("⏰", "Amazon EventBridge", ACCENT_VL,
         "Cron scheduler care declanseaza pipeline-ul automat zilnic la 06:00 UTC "
         "(08:00 ora Romaniei)."),
    ]

    y = Inches(1.85)
    for icon, name, color, role in services:
        h = Inches(0.55)
        # Icon column
        _add_card(s, Inches(0.5), y, Inches(0.7), h, fill=BG_SURFACE, border=color)
        _add_text(s, Inches(0.5), y + Inches(0.10), Inches(0.7), Inches(0.4),
                  icon, size=20, color=color, align=PP_ALIGN.CENTER)
        # Name column
        _add_card(s, Inches(1.30), y, Inches(2.7), h, fill=BG_SURFACE, border=color)
        _add_text(s, Inches(1.40), y + Inches(0.13), Inches(2.6), Inches(0.4),
                  name, size=13, bold=True, color=color)
        # Role column
        _add_card(s, Inches(4.10), y, Inches(8.7), h, fill=BG_DEEP, border=color)
        _add_text(s, Inches(4.30), y + Inches(0.13), Inches(8.5), Inches(0.4),
                  role, size=11, color=TEXT_LIGHT)
        y += h + Inches(0.04)


# ════════════════════════════════════════════════════════════════════
# SECTION 5 — DETAILED SERVICE SLIDES (7-15)
# ════════════════════════════════════════════════════════════════════

def _service_slide_skeleton(prs, *, title, subtitle, color,
                            about_text, role_text, body_builder):
    """Build a service slide with the standard layout:

      Title + subtitle
      ┌─DESPRE SERVICIU──────┐  ┌─ROL IN WEATHERLENS────┐
      │ {about_text}          │  │ {role_text}            │
      └───────────────────────┘  └────────────────────────┘
      {body content rendered by body_builder(slide)}
    """
    s = _new_blank(prs)
    _slide_title(s, title, subtitle, accent=color)
    _twin_callouts(
        s,
        left_label="DESPRE SERVICIU",
        left_text=about_text,
        right_label="ROL IN WEATHERLENS",
        right_text=role_text,
        right_color=color,
    )
    body_builder(s, color)
    return s


def slide_07_s3(prs):
    def body(s, color):
        buckets = [
            ("weatherlens-raw",       "JSON brut, partitionat year=Y/month=M/day=D/"),
            ("weatherlens-processed", "Parquet+Snappy: city=X/year=Y/month=M/data_*.parquet"),
            ("weatherlens-athena",    "Cache pentru rezultatele query-urilor Athena"),
            ("weatherlens-exports",   "Excel / CSV / JSON generate la cerere"),
            ("weatherlens-frontend",  "Static website hosting pentru dashboard HTML/JS"),
        ]
        y = Inches(3.45)
        for name, desc in buckets:
            h = Inches(0.50)
            _add_card(s, Inches(0.5), y, Inches(12.3), h, fill=BG_SURFACE, border=color)
            _add_text(s, Inches(0.7), y + Inches(0.10), Inches(4.5), Inches(0.4),
                      name, size=13, bold=True, color=color)
            _add_text(s, Inches(5.3), y + Inches(0.10), Inches(7.4), Inches(0.4),
                      desc, size=12, color=TEXT_LIGHT)
            y += h + Inches(0.08)
        _add_text(s, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4),
                  "Versioning activat · Lifecycle: raw > 90 zile auto-deleted · Snappy: 10x compresie",
                  size=11, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    _service_slide_skeleton(
        prs,
        title="Amazon S3 — Data Lake",
        subtitle="Storage layer · 5 bucket-uri partitionate Hive-style",
        color=ACCENT,
        about_text=(
            "S3 (Simple Storage Service) este serviciul AWS de stocare obiect — "
            "practic un hard disk infinit in cloud, accesibil prin API, cu cost "
            "~$0.023/GB/luna. Folosit ca data lake pentru orice format (JSON, "
            "Parquet, CSV, imagini)."
        ),
        role_text=(
            "Coloana vertebrala de storage a proiectului — 5 bucket-uri pentru date "
            "brute, procesate, cache Athena, exporturi user-facing si site-ul static "
            "al frontend-ului. Partitionarea Hive city/year/month reduce cu 90% "
            "datele scanate la query-uri."
        ),
        body_builder=body,
    )


def slide_08_lambda(prs):
    def body(s, color):
        fns = [
            ("weatherlens-ingest",   "Apeleaza Open-Meteo pentru 5 orase, salveaza JSON in S3 raw. Cron 06:00 UTC."),
            ("weatherlens-validate", "Verifica schema JSON inainte de ETL. Apelata din Step Functions."),
            ("weatherlens-export",   "Genereaza Excel/CSV/JSON la cerere. Ruleaza Athena, returneaza URL presemnat."),
        ]
        y = Inches(3.45)
        for name, desc in fns:
            h = Inches(0.95)
            _add_card(s, Inches(0.5), y, Inches(12.3), h, fill=BG_SURFACE, border=color)
            _add_text(s, Inches(0.7), y + Inches(0.13), Inches(11.9), Inches(0.4),
                      name, size=15, bold=True, color=color)
            _add_text(s, Inches(0.7), y + Inches(0.50), Inches(11.9), Inches(0.4),
                      desc, size=12, color=TEXT_LIGHT)
            y += h + Inches(0.12)
        _add_text(s, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4),
                  "Runtime: Python 3.11 · Memory: 256 MB · Timeout: 300s · Layer cu requests + boto3",
                  size=11, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    _service_slide_skeleton(
        prs,
        title="AWS Lambda — Compute serverless",
        subtitle="3 functii Python · pay-per-use · zero infra de mentenut",
        color=ACCENT_AM,
        about_text=(
            "Lambda este modelul Function-as-a-Service de la AWS — scrii o functie "
            "Python, AWS o ruleaza la cerere fara servere de configurat. Platesti "
            "doar pentru milisecundele in care ruleaza efectiv; Free Tier acopera "
            "1 milion de invocari/luna."
        ),
        role_text=(
            "Ruleaza cele 3 functii Python ale proiectului: ingest zilnic din Open-Meteo "
            "(cron-triggered), validare schema JSON inainte de ETL (sfn-triggered) "
            "si generare exports Excel/CSV/JSON la cerere din frontend "
            "(api-triggered)."
        ),
        body_builder=body,
    )


def slide_09_glue(prs):
    def body(s, color):
        items = [
            "Flatten JSON arrays paralele → 1 rand per ora",
            "Parse timestamp → date, hour, day_of_week, season, is_daytime",
            "Map weather_code (WMO) → weather_description + weather_category",
            "Validate ranges: temp ∈ [-50, 55], wind ∈ [0, 300] → null daca out-of-range",
            "Rolling avg 24h cu PySpark Window functions (per oras)",
            "Anomaly = temp curent - daily_avg per oras",
            "Wind category: Calm / Light / Moderate / Strong / Storm",
            "Extreme event flag: temp extrema, ploaie >20 mm, vant >80 km/h",
            "Dedupe pe (city, timestamp) si write Parquet+Snappy partitionat",
        ]
        _add_bullets(s, Inches(0.7), Inches(3.45), Inches(12), Inches(3.5),
                     items, size=12, space_after=3)

    _service_slide_skeleton(
        prs,
        title="AWS Glue — ETL serverless cu PySpark",
        subtitle="JSON → Parquet · 13 coloane derivate · validare automata",
        color=ACCENT_HOT,
        about_text=(
            "Glue este versiunea managed a Apache Spark — un cluster care apare la "
            "cerere, ruleaza job-ul ETL si se opreste automat. Avantaj cheie: "
            "scaleaza de la 100 MB la 100 TB fara modificari de cod, doar "
            "ajustand numarul de workers."
        ),
        role_text=(
            "Motorul ETL care citeste JSON brut din S3 raw, aplica 13 derivari "
            "analitice (sezon, rolling avg 24h, anomalii, categorii vant, flag "
            "fenomene extreme) plus validare ranges, si scrie Parquet partitionat. "
            "La volumul nostru ruleaza in <3 minute pe 2 workers G.1X."
        ),
        body_builder=body,
    )


def slide_10_step_functions(prs):
    def body(s, color):
        # Mini state machine diagram
        states = [("Validate", ACCENT), ("RunGlueETL", ACCENT_HOT),
                  ("RunCrawler", ACCENT_AM), ("EmitMetrics", ACCENT_GR)]
        box_w = Inches(2.5)
        box_h = Inches(0.6)
        total = len(states) * box_w + (len(states) - 1) * Inches(0.22)
        start_x = (WIDE_W - total) / 2
        y = Inches(3.45)
        for i, (name, c) in enumerate(states):
            x = start_x + i * (box_w + Inches(0.22))
            _add_card(s, x, y, box_w, box_h, fill=BG_SURFACE, border=c)
            _add_text(s, x, y + Inches(0.13), box_w, Inches(0.4),
                      name, size=13, bold=True, color=c, align=PP_ALIGN.CENTER)
            if i < len(states) - 1:
                _add_text(s, x + box_w, y + Inches(0.13), Inches(0.22), Inches(0.4),
                          "→", size=18, bold=True, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

        items = [
            "Vizualizare grafica live a executiei — vezi instant ce pas a esuat",
            "Retry automat cu backoff exponential — fara cod custom",
            "Catch / Throw declarativ pentru gestionarea erorilor",
            "Audit trail complet — log per executie pentru debugging",
            "Cost: ~$0.0002 per executie ⇒ practic gratuit",
        ]
        _add_bullets(s, Inches(0.7), Inches(4.30), Inches(12), Inches(2.7),
                     items, size=13, space_after=4)

    _service_slide_skeleton(
        prs,
        title="AWS Step Functions ⭐ — Serviciu nou (Nota 10)",
        subtitle="Orchestrare vizuala state machine cu retry/catch automat",
        color=ACCENT_GR,
        about_text=(
            "Step Functions defineste pipeline-ul ca un graf de stari in JSON "
            "(Amazon States Language) — fiecare pas are retry policy si error "
            "handler declarativi. Vizualizarea grafica live transforma "
            "debugging-ul dintr-o cautare prin loguri intr-o privire pe diagrama."
        ),
        role_text=(
            "Orchestreaza pipeline-ul ETL printr-un state machine cu 4 pasi: "
            "declanseaza Lambda de validare, asteapta Glue Job-ul sa se termine, "
            "ruleaza Glue Crawler-ul si emite metricile finale catre CloudWatch. "
            "Retry exponential la fiecare pas, automat."
        ),
        body_builder=body,
    )


def slide_11_athena(prs):
    def body(s, color):
        items = [
            "weatherlens_db.weather_hourly: 14 280 randuri orare (tabela principala)",
            "weatherlens_db.daily_summary: 595 rezumate zilnice (tabela aggregate)",
            "Partitionare city/year/month — Athena scaneaza doar partitiile cerute",
            "10 query-uri salvate ca Named Queries (reutilizabile)",
        ]
        _add_bullets(s, Inches(0.7), Inches(3.45), Inches(12), Inches(1.7),
                     items, size=13, space_after=4)

        _add_card(s, Inches(0.7), Inches(5.20), Inches(12), Inches(1.85),
                  fill=BG_SURFACE, border=color)
        code = (
            "SELECT date, city,\n"
            "       ROUND(avg_temp, 1) AS avg_temperature,\n"
            "       ROUND(total_precipitation, 1) AS rain_mm\n"
            "FROM weatherlens_db.daily_summary\n"
            "WHERE date >= CURRENT_DATE - INTERVAL '90' DAY\n"
            "ORDER BY city, date;"
        )
        box = s.shapes.add_textbox(Inches(0.9), Inches(5.32), Inches(11.6), Inches(1.6))
        tf = box.text_frame
        tf.word_wrap = True
        for i, line in enumerate(code.split("\n")):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = para.add_run()
            run.text = line
            run.font.size = Pt(11)
            run.font.name = "Consolas"
            run.font.color.rgb = TEXT_LIGHT

    _service_slide_skeleton(
        prs,
        title="Amazon Athena — SQL serverless",
        subtitle="Query-uri SQL pe S3 Parquet · pay-per-query · fara baza de date",
        color=ACCENT,
        about_text=(
            "Athena este Presto-as-a-Service de la AWS — un motor de query distribuit "
            "care intelege SQL ANSI dar citeste datele direct din fisiere "
            "Parquet/CSV/JSON din S3. Costul $5/TB scanat face partitionarea direct "
            "un instrument de optimizare a costurilor."
        ),
        role_text=(
            "Ruleaza cele 10 query-uri pre-definite ale proiectului direct pe "
            "Parquet din S3, fara baza de date in mijloc. Query-urile (trend "
            "temperatura, distributie precipitatii, top zile extreme, anomalii) "
            "se executa in 1-3 secunde si costa ~$0.00003 fiecare."
        ),
        body_builder=body,
    )


def slide_12_quicksight(prs):
    def body(s, color):
        charts = [
            ("Line chart",    "Evolutie temperatura zilnica per oras"),
            ("Bar chart",     "Precipitatii lunare sortate descendent"),
            ("Heat map",      "Calendar temperaturi"),
            ("Donut chart",   "Distributie tipuri de vreme"),
            ("Polar chart",   "Wind rose pe categorii vant"),
            ("KPI cards",     "Temp · Ploaie · Vant · Extreme"),
            ("Anomaly area",  "Temperatura vs medie istorica"),
        ]
        cell_w = Inches(2.95)
        cell_h = Inches(1.30)
        cols = 4
        start_x = Inches(0.5)
        start_y = Inches(3.45)
        for i, (kind, desc) in enumerate(charts):
            col = i % cols
            row = i // cols
            x = start_x + col * (cell_w + Inches(0.15))
            y = start_y + row * (cell_h + Inches(0.15))
            _add_card(s, x, y, cell_w, cell_h, fill=BG_SURFACE, border=color)
            _add_text(s, x + Inches(0.15), y + Inches(0.13), cell_w - Inches(0.3),
                      Inches(0.4), kind, size=13, bold=True, color=color)
            _add_text(s, x + Inches(0.15), y + Inches(0.5), cell_w - Inches(0.3),
                      Inches(0.75), desc, size=11, color=TEXT_LIGHT)

    _service_slide_skeleton(
        prs,
        title="Amazon QuickSight — Dashboard BI",
        subtitle="Vizualizari interactive · filtre · refresh automat din SPICE",
        color=ACCENT_HOT,
        about_text=(
            "QuickSight este alternativa AWS la Tableau / PowerBI — un BI tool web "
            "cu conexiune nativa la Athena, RDS sau direct la S3. Cache-ul SPICE "
            "precalculeaza rezultatele asa incat utilizatorul nu asteapta Athena la "
            "fiecare interactiune cu filtrele."
        ),
        role_text=(
            "Dashboard-ul BI nativ AWS al proiectului — preia date din Athena, "
            "afiseaza 7 grafice interactive (line, bar, donut, polar, heatmap, "
            "KPI cards), permite filtre dinamice (oras, perioada, sezon) si refresh "
            "automat zilnic dupa ce pipeline-ul ETL a terminat."
        ),
        body_builder=body,
    )


def slide_13_cloudwatch(prs):
    def body(s, color):
        items = [
            "Log Groups: /aws/lambda/weatherlens-* · /aws-glue/jobs/* · /aws/states/*",
            "Retention 30 zile · cautare full-text in CloudWatch Insights",
            "Custom metrics: PipelineSuccess, RecordsIngested, DurationSeconds, DataQualityScore",
            "Alarm: weatherlens-pipeline-failure → SNS email cand metric > 0",
            "CloudWatch Dashboard operational: success rate, durata medie, erori 24h",
        ]
        _add_bullets(s, Inches(0.7), Inches(3.45), Inches(12), Inches(3.0),
                     items, size=13, space_after=6)

    _service_slide_skeleton(
        prs,
        title="Amazon CloudWatch — Monitoring",
        subtitle="Logs centralizate · Metrici custom · Alarme",
        color=ACCENT_AM,
        about_text=(
            "CloudWatch este sistemul AWS de observabilitate — colecteaza orice log "
            "si metric din orice serviciu, cu interogare full-text si retentie "
            "configurabila. Alarms se leaga la SNS pentru notificari automate cand "
            "un metric depaseste pragul."
        ),
        role_text=(
            "Centralizeaza toate log-urile si metricile pipeline-ului — fiecare "
            "Lambda, Glue Job si executie Step Functions trimit acolo. Alarme "
            "email cand pipeline-ul esueaza, plus metrici custom (PipelineSuccess, "
            "RecordsIngested, DurationSeconds)."
        ),
        body_builder=body,
    )


def slide_14_iam(prs):
    def body(s, color):
        roles = [
            ("weatherlens-glue-role",
             "Trust: glue.amazonaws.com · Permite: S3 read/write, CloudWatch logs"),
            ("weatherlens-lambda-role",
             "Trust: lambda.amazonaws.com · Permite: S3, SFN invoke, Athena, CloudWatch"),
            ("weatherlens-sfn-role",
             "Trust: states.amazonaws.com · Permite: Lambda invoke, Glue, CloudWatch metrics"),
        ]
        y = Inches(3.45)
        for name, desc in roles:
            h = Inches(1.00)
            _add_card(s, Inches(0.5), y, Inches(12.3), h, fill=BG_SURFACE, border=color)
            _add_text(s, Inches(0.7), y + Inches(0.13), Inches(11.9), Inches(0.45),
                      name, size=15, bold=True, color=color)
            _add_text(s, Inches(0.7), y + Inches(0.52), Inches(11.9), Inches(0.5),
                      desc, size=12, color=TEXT_LIGHT)
            y += h + Inches(0.10)

    _service_slide_skeleton(
        prs,
        title="AWS IAM — Securitate fine-grained",
        subtitle="Roluri least-privilege · zero credentials hardcoded",
        color=ACCENT_GR,
        about_text=(
            "IAM (Identity and Access Management) este sistemul central de control "
            "al accesului in AWS — defineste cine (roluri, useri) poate face ce "
            "(actions) pe care resurse. Principiul least-privilege izoleaza serviciile "
            "intre ele si reduce raza de impact a unei compromisiuni."
        ),
        role_text=(
            "Defineste 3 roluri dedicate (pentru Glue, Lambda si Step Functions) cu "
            "policies fine-grained care permit doar operatiile minime necesare. "
            "Zero credentials hardcoded — toate apelurile folosesc role assumption "
            "automat."
        ),
        body_builder=body,
    )


def slide_15_eventbridge(prs):
    def body(s, color):
        items = [
            "Rule: weatherlens-daily-ingest cu schedule-expression cron(0 6 * * ? *)",
            "Target: arn:aws:lambda:eu-central-1:*:function:weatherlens-ingest",
            "State: ENABLED · ruleaza zilnic la 06:00 UTC = 08:00 ora Romaniei",
            "Free Tier: 1M events/luna ⇒ proiectul foloseste ~30 events/luna ⇒ $0",
            "Echivalente locale: APScheduler · Windows Task Scheduler · GitHub Actions cron",
        ]
        _add_bullets(s, Inches(0.7), Inches(3.45), Inches(12), Inches(3.5),
                     items, size=13, space_after=6)

    _service_slide_skeleton(
        prs,
        title="Amazon EventBridge — Cron scheduler",
        subtitle="Declansare zilnica automata · zero interventie manuala",
        color=ACCENT_VL,
        about_text=(
            "EventBridge este event bus-ul AWS — un sistem publish/subscribe in care "
            "orice eveniment (cron tick, S3 PutObject, schimbare de stare) poate "
            "triggera orice serviciu. La noi e folosit doar in rol de cron, dar "
            "puterea reala apare in pipelines reactive complete."
        ),
        role_text=(
            "Scheduler-ul cron care declanseaza pipeline-ul automat in fiecare zi la "
            "06:00 UTC — apeleaza Lambda-ul de ingest, care porneste reactia in lant: "
            "Lambda → S3 Event → Step Functions → Glue → Athena → refresh QuickSight."
        ),
        body_builder=body,
    )


# ════════════════════════════════════════════════════════════════════
# SECTION 6 — CONNECTION SCHEMA
# ════════════════════════════════════════════════════════════════════

def slide_16_connection_schema(prs):
    s = _new_blank(prs)
    _slide_title(s, "Schema conexiunilor — cum se leaga cele 9 servicii",
                 "Diagrama statica a relatiilor intre serviciile AWS folosite")

    # Box helper
    def _svc(x, y, icon, name, color, w=Inches(2.4), h=Inches(0.95)):
        _add_card(s, x, y, w, h, fill=BG_SURFACE, border=color, border_width=1.5)
        _add_text(s, x, y + Inches(0.10), w, Inches(0.35),
                  icon, size=18, color=color, align=PP_ALIGN.CENTER)
        _add_text(s, x, y + Inches(0.50), w, Inches(0.4),
                  name, size=11, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

    def _arrow(x, y, w=Inches(0.4), text="→"):
        _add_text(s, x, y, w, Inches(0.4),
                  text, size=20, bold=True, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    # ─── Top trigger row ─────────────────────────
    # ⏰ EventBridge (top-left center)
    _svc(Inches(0.6), Inches(1.85), "⏰", "EventBridge", ACCENT_VL)
    _arrow(Inches(3.0), Inches(2.18))
    # 🔧 Lambda ingest
    _svc(Inches(3.4), Inches(1.85), "🔧", "Lambda ingest", ACCENT_AM)
    _arrow(Inches(5.8), Inches(2.18))
    # 📡 Open-Meteo (external — dimmed)
    _add_card(s, Inches(6.2), Inches(1.85), Inches(2.4), Inches(0.95),
              fill=BG_DEEP, border=TEXT_DIM, border_width=1.0)
    _add_text(s, Inches(6.2), Inches(1.95), Inches(2.4), Inches(0.35),
              "📡", size=18, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    _add_text(s, Inches(6.2), Inches(2.35), Inches(2.4), Inches(0.4),
              "Open-Meteo API", size=11, bold=True, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    # Down arrow from Lambda ingest
    _add_text(s, Inches(4.4), Inches(2.85), Inches(0.4), Inches(0.4),
              "↓", size=20, bold=True, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    # ─── Storage layer (S3 raw) ──────────────────
    _svc(Inches(3.4), Inches(3.30), "📦", "S3 raw (JSON)", ACCENT)
    # Right arrow into Step Functions
    _arrow(Inches(5.8), Inches(3.63))
    _svc(Inches(6.2), Inches(3.30), "🎯", "Step Functions", ACCENT_GR)

    # Down arrow under Step Functions
    _add_text(s, Inches(7.2), Inches(4.30), Inches(0.4), Inches(0.4),
              "↓", size=20, bold=True, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    # SFN orchestrates Lambda validate + Glue ETL
    _svc(Inches(5.0), Inches(4.75), "🔍", "Lambda validate", ACCENT_AM,
         w=Inches(2.2), h=Inches(0.85))
    _add_text(s, Inches(7.2), Inches(4.95), Inches(0.4), Inches(0.4),
              "→", size=20, bold=True, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    _svc(Inches(7.6), Inches(4.75), "⚙️", "Glue ETL", ACCENT_HOT,
         w=Inches(2.2), h=Inches(0.85))
    _add_text(s, Inches(9.8), Inches(4.95), Inches(0.4), Inches(0.4),
              "→", size=20, bold=True, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    _svc(Inches(10.2), Inches(4.75), "📦", "S3 processed", ACCENT,
         w=Inches(2.6), h=Inches(0.85))

    # ─── Bottom analytics row ────────────────────
    # Down arrow from S3 processed
    _add_text(s, Inches(11.3), Inches(5.65), Inches(0.4), Inches(0.4),
              "↓", size=20, bold=True, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    _svc(Inches(0.6), Inches(6.10), "🗄️", "Athena", ACCENT,
         w=Inches(2.4), h=Inches(0.80))
    _add_text(s, Inches(3.0), Inches(6.30), Inches(0.4), Inches(0.4),
              "→", size=20, bold=True, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    _svc(Inches(3.4), Inches(6.10), "📊", "QuickSight", ACCENT_HOT,
         w=Inches(2.4), h=Inches(0.80))
    _add_text(s, Inches(5.8), Inches(6.30), Inches(0.4), Inches(0.4),
              "←", size=20, bold=True, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    _svc(Inches(6.2), Inches(6.10), "📥", "Lambda export", ACCENT_AM,
         w=Inches(2.4), h=Inches(0.80))

    # ─── Cross-cutting (right side) ──────────────
    cc_x = Inches(9.0)
    cc_y = Inches(6.10)
    _add_card(s, cc_x, cc_y, Inches(3.8), Inches(0.80),
              fill=BG_SURFACE, border=ACCENT_VL, border_width=1.0)
    _add_text(s, cc_x + Inches(0.15), cc_y + Inches(0.10), Inches(3.6), Inches(0.3),
              "CROSS-CUTTING", size=10, bold=True, color=ACCENT_VL)
    _add_text(s, cc_x + Inches(0.15), cc_y + Inches(0.40), Inches(3.6), Inches(0.4),
              "📈 CloudWatch · 🔐 IAM",
              size=12, bold=True, color=TEXT_LIGHT)


# ════════════════════════════════════════════════════════════════════
# SECTION 7 — MACRO WORKFLOW (4 phases)
# ════════════════════════════════════════════════════════════════════

def slide_17_macro_workflow(prs):
    s = _new_blank(prs)
    _slide_title(s, "Workflow la nivel macro — 4 faze",
                 "Cum curge data prin proiect, de la API la utilizator")

    phases = [
        ("1. PRELUARE",
         "Date externe → S3",
         ["EventBridge cron", "Lambda ingest", "Open-Meteo API", "S3 raw (JSON)"],
         ACCENT_AM),
        ("2. PRELUCRARE",
         "JSON → Parquet curatat",
         ["Step Functions", "Lambda validate", "Glue ETL (PySpark)", "13 derivari analitice"],
         ACCENT_HOT),
        ("3. STOCARE",
         "Catalog SQL-queryable",
         ["S3 processed (Parquet)", "Glue Data Catalog", "Athena views", "Partitionare Hive"],
         ACCENT),
        ("4. AFISARE",
         "Date → utilizator final",
         ["QuickSight dashboard", "Frontend HTML/JS", "Lambda export", "CSV / Excel / JSON"],
         ACCENT_GR),
    ]
    card_w = Inches(2.95)
    card_h = Inches(4.2)
    arrow_w = Inches(0.25)
    total_w = 4 * card_w + 3 * arrow_w
    start_x = (WIDE_W - total_w) / 2
    y = Inches(2.0)

    for i, (phase, subtitle, items, color) in enumerate(phases):
        x = start_x + i * (card_w + arrow_w)
        _add_card(s, x, y, card_w, card_h, fill=BG_SURFACE, border=color,
                  border_width=2.0)
        # Phase number + title
        _add_text(s, x, y + Inches(0.15), card_w, Inches(0.40),
                  phase, size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
        # Subtitle
        _add_text(s, x, y + Inches(0.55), card_w, Inches(0.5),
                  subtitle, size=11, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
        # Separator line
        sep = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 x + Inches(0.4), y + Inches(1.1),
                                 card_w - Inches(0.8), Inches(0.012))
        sep.fill.solid()
        sep.fill.fore_color.rgb = TEXT_DIM
        sep.line.fill.background()
        sep.shadow.inherit = False
        # Items
        item_box = s.shapes.add_textbox(x + Inches(0.2), y + Inches(1.25),
                                        card_w - Inches(0.4), card_h - Inches(1.4))
        tf = item_box.text_frame
        tf.word_wrap = True
        for j, item in enumerate(items):
            para = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            para.alignment = PP_ALIGN.LEFT
            para.space_after = Pt(10)
            run = para.add_run()
            run.text = f"▸ {item}"
            run.font.size = Pt(13)
            run.font.color.rgb = TEXT_LIGHT
            run.font.name = "Calibri"

        # Arrow between cards
        if i < 3:
            _add_text(s, x + card_w, y + Inches(1.95), arrow_w, Inches(0.4),
                      "→", size=22, bold=True, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    # Bottom timeline
    _add_card(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.7),
              fill=BG_SURFACE, border=ACCENT_VL)
    _add_text(s, Inches(0.7), Inches(6.66), Inches(12.0), Inches(0.5),
              "06:00 UTC declansare → 06:03 UTC date in catalog → la cerere user vede dashboard / exporta date",
              size=12, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SECTION 8 — CONCLUSION
# ════════════════════════════════════════════════════════════════════

def slide_18_conclusion(prs):
    s = _new_blank(prs)
    _slide_title(s, "Concluzii",
                 "9 servicii AWS · pipeline complet · arhitectura production-ready")

    items = [
        "Pipeline ETL serverless complet integrat cu 9 servicii AWS — depaseste cerinta de 8 servicii",
        "Step Functions ⭐ ca serviciu nou — orchestrare vizuala cu retry / catch automat",
        "Acoperire layere: Trigger · Ingestion · Orchestration · Processing · Analytics · Visualization",
        "Dataset real: 5 orase romanesti, 119 zile, 14 280 randuri orare procesate",
        "Cost AWS estimat: ~$2–20/luna · cost local-first: $0",
        "12 teste pytest verzi · documentatie completa · 4 pagini dashboard livrate",
    ]
    _add_bullets(s, Inches(0.7), Inches(1.95), Inches(12), Inches(3.5),
                 items, size=15, space_after=8)

    # Closing card
    _add_card(s, Inches(0.7), Inches(5.6), Inches(12), Inches(1.4),
              fill=BG_SURFACE, border=ACCENT_GR)
    _add_text(s, Inches(0.9), Inches(5.78), Inches(11.6), Inches(0.6),
              "Multumesc!",
              size=32, bold=True, color=ACCENT_GR, align=PP_ALIGN.CENTER)
    _add_text(s, Inches(0.9), Inches(6.40), Inches(11.6), Inches(0.5),
              "Demo: streamlit run dashboard/app.py",
              size=14, color=TEXT_MUTED, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# MAIN
# ════════════════════════════════════════════════════════════════════

def build() -> Path:
    prs = Presentation()
    prs.slide_width = WIDE_W
    prs.slide_height = WIDE_H

    builders = [
        slide_01_title,
        slide_02_what_is_aws,
        slide_03_aws_categories,
        slide_04_weatherlens_app,
        slide_05_scope,
        slide_06_services_list,
        slide_07_s3,
        slide_08_lambda,
        slide_09_glue,
        slide_10_step_functions,
        slide_11_athena,
        slide_12_quicksight,
        slide_13_cloudwatch,
        slide_14_iam,
        slide_15_eventbridge,
        slide_16_connection_schema,
        slide_17_macro_workflow,
        slide_18_conclusion,
    ]
    for b in builders:
        b(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    path = build()
    print(f"Saved: {path}")
    print(f"Slides: 18")
